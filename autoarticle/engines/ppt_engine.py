
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import json
import datetime
from io import BytesIO

from .constants import THEMES
from .utils import get_valid_images

class PPTEngine:
    def __init__(self, theme_name, school_name, layout_version="v1"):
        self.theme_name = theme_name
        self.theme = THEMES.get(theme_name, THEMES["웜 & 플레이풀"])
        self.school_name = school_name
        self.layout_version = layout_version
        self.prs = Presentation()
        # Set slide size to 16:9 (13.333 x 7.5 inches)
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def create_presentation(self, articles, output_path=None):
        """
        Creates a PPT presentation from a list of articles.
        Returns bytes if output_path is None.
        """
        self._add_title_slide()
        
        for art in articles:
            self._add_article_slide(art)
            
        if output_path:
            self.prs.save(output_path)
            return output_path
        else:
            buffer = BytesIO()
            self.prs.save(buffer)
            buffer.seek(0)
            return buffer

    def _add_title_slide(self):
        if self.layout_version == "v6":
            return self._add_title_slide_v6()
        if self.layout_version == "v5":
            return self._add_title_slide_v5()
        if self.layout_version == "v4":
            return self._add_title_slide_v4()
        if self.layout_version == "v3":
            return self._add_title_slide_v3()
        if self.layout_version == "v2":
            return self._add_title_slide_v2()
        """Adds a professional main title slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6]) # Blank layout
        
        # Background Color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*self.theme["sub"])
        
        # Center Box for Title
        left = Inches(2)
        top = Inches(2.5)
        width = Inches(9.333)
        height = Inches(2.5)
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.line.width = Pt(0)
        shape.shadow.inherit = False
        
        # Title Text
        text_frame = shape.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.text = f"{self.school_name} 뉴스레터"
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Malgun Gothic'
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*self.theme["main"])
        
        # Subtitle (Date)
        p = text_frame.add_paragraph()
        now = datetime.datetime.now()
        p.text = f"{now.year}학년도 {now.month}월 소식"
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Malgun Gothic'
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(80, 80, 80)

    def _add_title_slide_v2(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(248, 249, 251)

        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.35), self.prs.slide_height)
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(*self.theme["main"])
        accent.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11.5), Inches(2.2))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = self.school_name
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(36, 36, 36)

        p2 = tf.add_paragraph()
        now = datetime.datetime.now()
        p2.text = f"{now.year}학년도 {now.month}월 뉴스레터"
        p2.font.name = "Malgun Gothic"
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(*self.theme["main"])

    def _add_title_slide_v3(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(242, 239, 232)

        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), self.prs.slide_width, Inches(1.1))
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(26, 30, 40)
        header.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(12), Inches(0.6))
        p = tb.text_frame.paragraphs[0]
        p.text = "SCHOOL MAGAZINE"
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        title = slide.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(12), Inches(1.4))
        t = title.text_frame.paragraphs[0]
        t.text = self.school_name
        t.font.name = "Malgun Gothic"
        t.font.size = Pt(56)
        t.font.bold = True
        t.font.color.rgb = RGBColor(20, 20, 20)

        now = datetime.datetime.now()
        sub = title.text_frame.add_paragraph()
        sub.text = f"{now.year}학년도 {now.month}월 뉴스레터"
        sub.font.name = "Malgun Gothic"
        sub.font.size = Pt(24)
        sub.font.color.rgb = RGBColor(*self.theme["main"])

    def _add_title_slide_v4(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(247, 247, 247)

        masthead = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), self.prs.slide_width, Inches(0.95))
        masthead.fill.solid()
        masthead.fill.fore_color.rgb = RGBColor(20, 20, 20)
        masthead.line.fill.background()

        head = masthead.text_frame.paragraphs[0]
        head.text = "SCHOOL NEWSPAPER"
        head.font.name = "Malgun Gothic"
        head.font.size = Pt(20)
        head.font.bold = True
        head.font.color.rgb = RGBColor(255, 255, 255)

        title = slide.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(11.8), Inches(1.5))
        p = title.text_frame.paragraphs[0]
        p.text = self.school_name
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = RGBColor(24, 24, 24)

        now = datetime.datetime.now()
        sub = title.text_frame.add_paragraph()
        sub.text = now.strftime("%Y.%m")
        sub.font.name = "Malgun Gothic"
        sub.font.size = Pt(24)
        sub.font.color.rgb = RGBColor(85, 85, 85)

    def _add_title_slide_v5(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(243, 239, 232)

        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.35), self.prs.slide_width, Inches(1.2))
        band.fill.solid()
        band.fill.fore_color.rgb = RGBColor(*self.theme["main"])
        band.line.fill.background()

        title = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(12), Inches(1.6))
        p = title.text_frame.paragraphs[0]
        p.text = self.school_name
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(28, 28, 28)

        now = datetime.datetime.now()
        s = title.text_frame.add_paragraph()
        s.text = f"{now.year}.{now.month} Monthly Story"
        s.font.name = "Malgun Gothic"
        s.font.size = Pt(22)
        s.font.color.rgb = RGBColor(*self.theme["main"])

    def _add_title_slide_v6(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(255, 255, 255)

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.5), self.prs.slide_width, Inches(0.06))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(*self.theme["main"])
        line.line.fill.background()

        title = slide.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(11.2), Inches(1.4))
        t = title.text_frame.paragraphs[0]
        t.text = self.school_name
        t.font.name = "Malgun Gothic"
        t.font.size = Pt(58)
        t.font.bold = True
        t.font.color.rgb = RGBColor(20, 20, 20)

        now = datetime.datetime.now()
        sub = title.text_frame.add_paragraph()
        sub.text = now.strftime("%Y-%m")
        sub.font.name = "Malgun Gothic"
        sub.font.size = Pt(20)
        sub.font.color.rgb = RGBColor(110, 110, 110)

    def _add_article_slide(self, article):
        if self.layout_version == "v6":
            return self._add_article_slide_v6(article)
        if self.layout_version == "v5":
            return self._add_article_slide_v5(article)
        if self.layout_version == "v4":
            return self._add_article_slide_v4(article)
        if self.layout_version == "v3":
            return self._add_article_slide_v3(article)
        if self.layout_version == "v2":
            return self._add_article_slide_v2(article)
        """Adds a content slide with 2-column layout (Text Left, Image Right)."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6]) # Blank layout
        
        # 1. Header Bar (Visual Identity)
        header_height = Inches(1.0)
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, header_height
        )
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(*self.theme["main"])
        header.line.fill.background() # No outline
        
        # Header Text (Article Title)
        tf = header.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = article.get('title', '제목 없음')
        p.font.name = 'Malgun Gothic'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT
        tf.margin_left = Inches(0.5)
        
        # 2. Layout Grid
        # Left Column (Text): 0.5" ~ 7.0" (Width 6.5")
        # Right Column (Image): 7.5" ~ 12.8" (Width 5.3")
        # Top: 1.5"
        
        # --- Left Column: Content ---
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.5), Inches(5.5))
        tf = left_box.text_frame
        tf.word_wrap = True
        
        # Metadata (Date/Loc)
        p = tf.add_paragraph()
        info_parts = []
        if article.get('date'): info_parts.append(f"{article['date']}")
        if article.get('location'): info_parts.append(f"{article['location']}")
        if article.get('grade'): info_parts.append(f"{article['grade']}")
        p.text = " | ".join(info_parts)
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(128, 128, 128)
        p.space_after = Pt(10)
        
        # Body Content
        content_text = article.get('content', '')
        # Basic clean up if summarized
        if isinstance(content_text, list): # AI summary case
            for item in content_text:
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(20)
                p.space_after = Pt(10)
                p.level = 0
        else: # Raw text case
            # Truncate long text for PPT readability
            p = tf.add_paragraph()
            # Check for summary failure message or overly long text
            if "요약에 실패" in content_text or len(content_text) > 300:
                truncated = content_text[:300].rsplit(' ', 1)[0] if len(content_text) > 300 else content_text
                p.text = truncated + "..." if len(content_text) > 300 else content_text
            else:
                p.text = content_text
            p.font.size = Pt(18)
            p.font.name = "Malgun Gothic"
            
        # --- Right Column: Images (1~4) ---
        imgs_raw = article.get('images', [])
        valid_imgs = get_valid_images(imgs_raw, max_count=4)
        self._add_image_grid(slide, valid_imgs, left=7.2, top=1.5, width=5.5, height=5.0)

    def _add_image_grid(self, slide, images, left, top, width, height):
        if not images:
            return

        # (l, t, w, h) in inches
        gap = 0.08
        count = len(images)

        if count == 1:
            rects = [(left, top, width, height)]
        elif count == 2:
            half_w = (width - gap) / 2
            rects = [
                (left, top, half_w, height),
                (left + half_w + gap, top, half_w, height),
            ]
        elif count == 3:
            big_w = (width * 0.64) - (gap / 2)
            small_w = width - big_w - gap
            half_h = (height - gap) / 2
            rects = [
                (left, top, big_w, height),
                (left + big_w + gap, top, small_w, half_h),
                (left + big_w + gap, top + half_h + gap, small_w, half_h),
            ]
        else:
            half_w = (width - gap) / 2
            half_h = (height - gap) / 2
            rects = [
                (left, top, half_w, half_h),
                (left + half_w + gap, top, half_w, half_h),
                (left, top + half_h + gap, half_w, half_h),
                (left + half_w + gap, top + half_h + gap, half_w, half_h),
            ]

        for img_data, (l, t, w, h) in zip(images, rects):
            try:
                pic = slide.shapes.add_picture(img_data, Inches(l), Inches(t))
                aspect = pic.width.inches / pic.height.inches
                target_aspect = w / h
                if aspect > target_aspect:
                    pic.height = Inches(h)
                    pic.width = Inches(h * aspect)
                else:
                    pic.width = Inches(w)
                    pic.height = Inches(w / aspect)

                pic.left = Inches(l + max((w - pic.width.inches) / 2, 0))
                pic.top = Inches(t + max((h - pic.height.inches) / 2, 0))
                pic.line.color.rgb = RGBColor(200, 200, 200)
                pic.line.width = Pt(1)
            except Exception:
                continue

    def _add_article_slide_v2(self, article):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(250, 250, 250)

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.8), Inches(1.1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = article.get("title", "제목 없음")
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = RGBColor(28, 28, 28)

        meta_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.45), Inches(6.2), Inches(0.5))
        mt = meta_box.text_frame
        info_parts = []
        if article.get("date"):
            info_parts.append(article["date"])
        if article.get("location"):
            info_parts.append(article["location"])
        if article.get("grade"):
            info_parts.append(article["grade"])
        m = mt.paragraphs[0]
        m.text = "  |  ".join(info_parts)
        m.font.name = "Malgun Gothic"
        m.font.size = Pt(13)
        m.font.color.rgb = RGBColor(110, 110, 110)

        left_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(6.2), Inches(4.9))
        tf = left_box.text_frame
        tf.word_wrap = True
        content_text = article.get("content", "")
        if isinstance(content_text, list):
            for item in content_text:
                cp = tf.add_paragraph()
                cp.text = f"• {item}"
                cp.font.name = "Malgun Gothic"
                cp.font.size = Pt(20)
                cp.space_after = Pt(8)
        else:
            cp = tf.add_paragraph()
            cp.text = content_text[:320] + ("..." if len(content_text) > 320 else "")
            cp.font.name = "Malgun Gothic"
            cp.font.size = Pt(18)

        imgs_raw = article.get("images", [])
        valid_imgs = get_valid_images(imgs_raw, max_count=4)
        self._add_image_grid(slide, valid_imgs, left=7.3, top=2.0, width=5.1, height=4.9)

    def _add_article_slide_v3(self, article):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(246, 244, 239)

        top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), self.prs.slide_width, Inches(0.9))
        top.fill.solid()
        top.fill.fore_color.rgb = RGBColor(26, 30, 40)
        top.line.fill.background()

        head = top.text_frame.paragraphs[0]
        head.text = article.get("title", "제목 없음")
        head.font.name = "Malgun Gothic"
        head.font.size = Pt(24)
        head.font.bold = True
        head.font.color.rgb = RGBColor(255, 255, 255)

        meta_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(12), Inches(0.5))
        info = []
        if article.get("date"):
            info.append(article["date"])
        if article.get("location"):
            info.append(article["location"])
        if article.get("grade"):
            info.append(article["grade"])
        mp = meta_box.text_frame.paragraphs[0]
        mp.text = " | ".join(info)
        mp.font.name = "Malgun Gothic"
        mp.font.size = Pt(13)
        mp.font.color.rgb = RGBColor(90, 90, 90)

        left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.1), Inches(5.3))
        tf = left_box.text_frame
        tf.word_wrap = True
        content_text = article.get("content", "")
        if isinstance(content_text, list):
            for item in content_text:
                p = tf.add_paragraph()
                p.text = f"- {item}"
                p.font.name = "Malgun Gothic"
                p.font.size = Pt(20)
                p.space_after = Pt(10)
        else:
            p = tf.add_paragraph()
            p.text = content_text[:320] + ("..." if len(content_text) > 320 else "")
            p.font.name = "Malgun Gothic"
            p.font.size = Pt(18)

        imgs_raw = article.get("images", [])
        valid_imgs = get_valid_images(imgs_raw, max_count=4)
        self._add_image_grid(slide, valid_imgs, left=7.2, top=1.8, width=5.3, height=5.3)

    def _add_article_slide_v4(self, article):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(247, 247, 247)

        top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), self.prs.slide_width, Inches(0.75))
        top.fill.solid()
        top.fill.fore_color.rgb = RGBColor(20, 20, 20)
        top.line.fill.background()

        head = top.text_frame.paragraphs[0]
        head.text = article.get("title", "Untitled")
        head.font.name = "Malgun Gothic"
        head.font.size = Pt(22)
        head.font.bold = True
        head.font.color.rgb = RGBColor(255, 255, 255)

        meta = slide.shapes.add_textbox(Inches(0.8), Inches(0.95), Inches(12), Inches(0.45))
        meta_parts = [p for p in [article.get("date"), article.get("location"), article.get("grade")] if p]
        mp = meta.text_frame.paragraphs[0]
        mp.text = " | ".join(meta_parts)
        mp.font.name = "Malgun Gothic"
        mp.font.size = Pt(13)
        mp.font.color.rgb = RGBColor(95, 95, 95)

        body = slide.shapes.add_textbox(Inches(0.8), Inches(1.55), Inches(6.0), Inches(5.6))
        bt = body.text_frame
        bt.word_wrap = True
        content_text = article.get("content", "")
        if isinstance(content_text, list):
            for item in content_text:
                p = bt.add_paragraph()
                p.text = f"- {item}"
                p.font.name = "Malgun Gothic"
                p.font.size = Pt(18)
                p.space_after = Pt(8)
        else:
            p = bt.add_paragraph()
            p.text = content_text[:360] + ("..." if len(content_text) > 360 else "")
            p.font.name = "Malgun Gothic"
            p.font.size = Pt(17)

        imgs_raw = article.get("images", [])
        valid_imgs = get_valid_images(imgs_raw, max_count=4)
        self._add_image_grid(slide, valid_imgs, left=7.0, top=1.55, width=5.5, height=5.6)

    def _add_article_slide_v5(self, article):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(243, 239, 232)

        title_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.55), Inches(11.8), Inches(1.1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = article.get("title", "Untitled")
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(26, 26, 26)

        info = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.6), Inches(11.8), Inches(0.5))
        info.fill.solid()
        info.fill.fore_color.rgb = RGBColor(255, 255, 255)
        info.line.color.rgb = RGBColor(218, 218, 218)
        info.line.width = Pt(1)
        ip = info.text_frame.paragraphs[0]
        ip.text = " | ".join([p for p in [article.get("date"), article.get("location"), article.get("grade")] if p])
        ip.font.name = "Malgun Gothic"
        ip.font.size = Pt(13)
        ip.font.color.rgb = RGBColor(95, 95, 95)

        body = slide.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(5.7), Inches(4.8))
        bt = body.text_frame
        bt.word_wrap = True
        content_text = article.get("content", "")
        if isinstance(content_text, list):
            for item in content_text:
                p = bt.add_paragraph()
                p.text = f"- {item}"
                p.font.name = "Malgun Gothic"
                p.font.size = Pt(18)
                p.space_after = Pt(8)
        else:
            p = bt.add_paragraph()
            p.text = content_text[:340] + ("..." if len(content_text) > 340 else "")
            p.font.name = "Malgun Gothic"
            p.font.size = Pt(17)

        imgs_raw = article.get("images", [])
        valid_imgs = get_valid_images(imgs_raw, max_count=4)
        self._add_image_grid(slide, valid_imgs, left=6.85, top=2.3, width=5.65, height=4.8)

    def _add_article_slide_v6(self, article):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(255, 255, 255)

        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.35), self.prs.slide_width, Inches(0.05))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*self.theme["main"])
        bar.line.fill.background()

        title = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(12), Inches(1.0))
        tp = title.text_frame.paragraphs[0]
        tp.text = article.get("title", "Untitled")
        tp.font.name = "Malgun Gothic"
        tp.font.size = Pt(34)
        tp.font.bold = True
        tp.font.color.rgb = RGBColor(20, 20, 20)

        meta = slide.shapes.add_textbox(Inches(0.8), Inches(1.72), Inches(12), Inches(0.45))
        mp = meta.text_frame.paragraphs[0]
        mp.text = " | ".join([p for p in [article.get("date"), article.get("location"), article.get("grade")] if p])
        mp.font.name = "Malgun Gothic"
        mp.font.size = Pt(12)
        mp.font.color.rgb = RGBColor(120, 120, 120)

        imgs_raw = article.get("images", [])
        valid_imgs = get_valid_images(imgs_raw, max_count=4)
        self._add_image_grid(slide, valid_imgs, left=0.8, top=2.2, width=6.0, height=4.8)

        body = slide.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(5.4), Inches(4.8))
        bt = body.text_frame
        bt.word_wrap = True
        content_text = article.get("content", "")
        if isinstance(content_text, list):
            for item in content_text:
                p = bt.add_paragraph()
                p.text = f"- {item}"
                p.font.name = "Malgun Gothic"
                p.font.size = Pt(17)
                p.space_after = Pt(7)
        else:
            p = bt.add_paragraph()
            p.text = content_text[:320] + ("..." if len(content_text) > 320 else "")
            p.font.name = "Malgun Gothic"
            p.font.size = Pt(16)
